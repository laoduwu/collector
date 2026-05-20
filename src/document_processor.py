"""文档处理器：PDF/DOCX/PPTX/TXT/MD 提取、摘要、双语翻译、HTML 组装

PDF 使用 Gemini Vision 逐批渲染页面并转换为结构化 HTML，
可正确还原表格、标题层级、加粗、彩色框、双栏等复杂排版。
"""
import base64
import json
import os
import time
import urllib.request
import urllib.error
from html.parser import HTMLParser
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

import requests

from utils.logger import logger
from video_processor import _llm_json, _translate_paragraphs


# ─── 辅助 ────────────────────────────────────────────────────

def _esc(text: str) -> str:
    return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


def _html_to_plain(html: str) -> str:
    """从 HTML 提取纯文字（供 AI 摘要/翻译）"""
    class _Ex(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.parts: List[str] = []
        def handle_data(self, data: str) -> None:
            self.parts.append(data)
    ex = _Ex()
    ex.feed(html)
    return ' '.join(ex.parts)


def _strip_code_fence(text: str) -> str:
    """去掉 Gemini 可能返回的 ```html ... ``` 包裹"""
    text = text.strip()
    if text.startswith('```'):
        lines = text.splitlines()
        # 去掉首行 (```html 或 ```)
        lines = lines[1:]
        # 去掉末行 (```)
        if lines and lines[-1].strip() == '```':
            lines = lines[:-1]
        text = '\n'.join(lines)
    return text.strip()


# ─── Storage 下载 ───────────────────────────────────────────

def _download_from_storage(storage_path: str) -> bytes:
    supabase_url = os.environ['SUPABASE_URL']
    service_key = os.environ['SUPABASE_SERVICE_ROLE_KEY']
    url = f'{supabase_url}/storage/v1/object/media-temp/{storage_path}'
    resp = requests.get(url, headers={'Authorization': f'Bearer {service_key}'}, timeout=60)
    resp.raise_for_status()
    return resp.content


# ─── Gemini Vision 调用 ──────────────────────────────────────

def _call_gemini_vision(prompt: str, images: List[Tuple[str, str]], timeout: int = 180) -> str:
    """调用 Gemini Vision，返回文本响应。images = [(mime_type, base64_data), ...]"""
    api_key = os.environ.get('GEMINI_API_KEY', '')
    if not api_key:
        raise RuntimeError('GEMINI_API_KEY 未设置')

    url = (
        'https://generativelanguage.googleapis.com/v1beta/models/'
        f'gemini-2.5-flash:generateContent?key={api_key}'
    )

    parts: List[Dict] = []
    for mime_type, b64_data in images:
        parts.append({'inline_data': {'mime_type': mime_type, 'data': b64_data}})
    parts.append({'text': prompt})

    body = json.dumps({'contents': [{'parts': parts}]}).encode()
    req = urllib.request.Request(
        url, data=body, headers={'Content-Type': 'application/json'}
    )

    for attempt in range(3):
        try:
            with urllib.request.urlopen(req, timeout=timeout) as resp:
                result = json.loads(resp.read())
            return result['candidates'][0]['content']['parts'][0]['text']
        except urllib.error.HTTPError as e:
            if e.code == 429 and attempt < 2:
                wait = 20 * (attempt + 1)
                logger.warning(f'Gemini Vision 429，{wait}s 后重试')
                time.sleep(wait)
            else:
                raise


_PDF_VISION_PROMPT = """\
以下是PDF文档的连续页面截图，请将这些页面的内容**完整**转换为HTML格式。

转换规则：
1. 文字内容必须完整保留，不得遗漏任何段落
2. 表格 → <table><thead><tr><th></th></tr></thead><tbody><tr><td></td></tr></tbody></table>
3. 标题根据视觉大小 → <h1>/<h2>/<h3>/<h4>
4. 加粗 → <b>，斜体 → <i>
5. 重点框/提示框/callout → <div class="callout" style="border:1px solid #d1d5db;\
padding:12px 16px;margin:12px 0;border-radius:6px;background:#f9fafb;">
6. 绿色边框或"推荐"类框 → style 加 border-left:4px solid #22c55e
7. 红色边框或"不推荐"/"警告"类框 → style 加 border-left:4px solid #ef4444
8. 其他有色边框框 → style 加 border-left:4px solid #f59e0b
9. 无序列表 → <ul><li>，有序列表 → <ol><li>
10. 普通段落 → <p>
11. 忽略页眉/页脚/页码
12. 多栏布局按阅读顺序线性输出
13. 不包含 <html>/<head>/<body> 标签，只输出正文 HTML 片段
14. 不要用 ```html``` 包裹
"""


# ─── PDF 提取（Vision 主路径 + dict 降级） ───────────────────

def _extract_pdf_vision(raw: bytes) -> Tuple[str, str, str]:
    """用 Gemini Vision 将 PDF 每页转为结构化 HTML。
    返回 (body_html, plain_text, title)
    """
    import fitz

    doc = fitz.open(stream=raw, filetype='pdf')
    title = (doc.metadata or {}).get('title', '') or ''
    total_pages = doc.page_count
    logger.info(f'PDF Vision 提取：共 {total_pages} 页')

    BATCH_SIZE = 5   # 每次 API 调用处理的页数
    DPI = 96         # 分辨率：794×1123px（A4），平衡质量与体积

    html_parts: List[str] = []
    plain_parts: List[str] = []

    pages = list(doc)
    for batch_start in range(0, len(pages), BATCH_SIZE):
        batch = pages[batch_start: batch_start + BATCH_SIZE]
        end = min(batch_start + BATCH_SIZE, total_pages)
        logger.info(f'  处理第 {batch_start + 1}–{end} 页（共 {total_pages} 页）')

        images: List[Tuple[str, str]] = []
        for page in batch:
            pix = page.get_pixmap(dpi=DPI)
            png_bytes = pix.tobytes('png')
            images.append(('image/png', base64.b64encode(png_bytes).decode()))

        try:
            chunk_html = _call_gemini_vision(_PDF_VISION_PROMPT, images, timeout=180)
            chunk_html = _strip_code_fence(chunk_html)
            if chunk_html:
                html_parts.append(chunk_html)
                plain_parts.append(_html_to_plain(chunk_html))
        except Exception as e:
            logger.warning(f'Vision 处理第 {batch_start + 1}–{end} 页失败，降级文字提取: {e}')
            for page in batch:
                text = page.get_text().strip()
                if text:
                    html_parts.append(f'<p>{_esc(text)}</p>')
                    plain_parts.append(text)

    doc.close()
    return '\n'.join(html_parts), '\n\n'.join(plain_parts), title


def _extract_pdf_dict(raw: bytes) -> Tuple[str, str, str]:
    """基于字体大小的结构化提取（降级备用）。返回 (body_html, plain_text, title)"""
    import fitz
    import statistics

    doc = fitz.open(stream=raw, filetype='pdf')
    title = (doc.metadata or {}).get('title', '') or ''

    all_sizes: List[float] = []
    for page in doc:
        for block in page.get_text('dict')['blocks']:
            if block.get('type') != 0:
                continue
            for line in block['lines']:
                for span in line['spans']:
                    if span['text'].strip():
                        all_sizes.append(span['size'])

    body_size = statistics.median(all_sizes) if all_sizes else 10.0
    h1_min, h2_min, h3_min = body_size * 1.7, body_size * 1.35, body_size * 1.15

    toc = doc.get_toc()
    toc_html = ''
    if toc:
        items = ['<div class="toc"><h2>目录</h2><ul>']
        for level, t, _ in toc[:80]:
            pad = '　' * (level - 1)
            items.append(f'<li>{pad}{_esc(t)}</li>')
        toc_html = ''.join(items) + '</ul></div><hr>\n'

    html_parts = [toc_html]
    plain_parts: List[str] = []
    seen_xrefs: set = set()

    for page in doc:
        page_dict = page.get_text('dict')
        xref_to_uri: Dict[int, str] = {}
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            if xref in seen_xrefs:
                continue
            try:
                img_data = doc.extract_image(xref)
                if not img_data:
                    continue
                ext = img_data.get('ext', 'png')
                b64 = base64.b64encode(img_data['image']).decode()
                xref_to_uri[xref] = f'data:image/{ext};base64,{b64}'
                seen_xrefs.add(xref)
            except Exception:
                pass

        for block in page_dict['blocks']:
            btype = block.get('type')
            if btype == 1:
                xref = block.get('xref', 0)
                uri = xref_to_uri.get(xref)
                if uri:
                    html_parts.append(f'<p><img src="{uri}" /></p>')
                continue
            if btype != 0:
                continue

            block_html_lines: List[str] = []
            block_plain_lines: List[str] = []
            block_sizes: List[float] = []

            for line in block['lines']:
                span_html: List[str] = []
                span_plain: List[str] = []
                for span in line['spans']:
                    text = span['text']
                    if not text.strip():
                        continue
                    size = span['size']
                    flags = span.get('flags', 0)
                    block_sizes.append(size)
                    t = _esc(text)
                    if bool(flags & 8):
                        t = f'<code>{t}</code>'
                    else:
                        if bool(flags & 16):
                            t = f'<b>{t}</b>'
                        if bool(flags & 2):
                            t = f'<i>{t}</i>'
                    span_html.append(t)
                    span_plain.append(text)
                if span_html:
                    block_html_lines.append(''.join(span_html))
                    block_plain_lines.append(''.join(span_plain))

            if not block_html_lines:
                continue

            combined_html = ' '.join(block_html_lines)
            combined_plain = ' '.join(block_plain_lines)
            avg_size = statistics.mean(block_sizes) if block_sizes else body_size

            if avg_size >= h1_min:
                html_parts.append(f'<h1>{combined_html}</h1>')
            elif avg_size >= h2_min:
                html_parts.append(f'<h2>{combined_html}</h2>')
            elif avg_size >= h3_min:
                html_parts.append(f'<h3>{combined_html}</h3>')
            else:
                html_parts.append(f'<p>{combined_html}</p>')
            plain_parts.append(combined_plain)

    doc.close()
    return '\n'.join(html_parts), '\n\n'.join(plain_parts), title


# ─── DOCX 提取 ──────────────────────────────────────────────

def _extract_docx(raw: bytes) -> Tuple[str, str, str]:
    from docx import Document
    doc = Document(BytesIO(raw))
    title = (doc.core_properties.title or '').strip()
    html_parts: List[str] = []
    plain_parts: List[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = para.style.name if para.style else ''
        plain_parts.append(text)

        if 'Heading 1' in style_name:
            html_parts.append(f'<h1>{_esc(text)}</h1>')
        elif 'Heading 2' in style_name:
            html_parts.append(f'<h2>{_esc(text)}</h2>')
        elif 'Heading 3' in style_name or 'Heading 4' in style_name:
            html_parts.append(f'<h3>{_esc(text)}</h3>')
        else:
            run_parts: List[str] = []
            for run in para.runs:
                if not run.text:
                    continue
                t = _esc(run.text)
                if run.bold:
                    t = f'<b>{t}</b>'
                if run.italic:
                    t = f'<i>{t}</i>'
                run_parts.append(t)
            if run_parts:
                html_parts.append('<p>' + ''.join(run_parts) + '</p>')

    return '\n'.join(html_parts), '\n\n'.join(plain_parts), title


# ─── PPTX 提取 ──────────────────────────────────────────────

def _extract_pptx(raw: bytes) -> Tuple[str, str, str]:
    from pptx import Presentation
    prs = Presentation(BytesIO(raw))
    title = (prs.core_properties.title or '').strip()
    html_parts: List[str] = []
    plain_parts: List[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            html_parts.append(f'<h2>幻灯片 {i}</h2>')
            for t in slide_texts:
                html_parts.append(f'<p>{_esc(t)}</p>')
            plain_parts.append(f'[幻灯片 {i}]\n' + '\n'.join(slide_texts))

    return '\n'.join(html_parts), '\n\n'.join(plain_parts), title


# ─── TXT/MD 提取 ─────────────────────────────────────────────

def _extract_text_file(raw: bytes) -> Tuple[str, str, str]:
    text = raw.decode('utf-8', errors='replace')
    first_line = text.split('\n')[0].strip().lstrip('#').strip()
    title = first_line[:80] if first_line else ''
    paras = [p.strip() for p in text.split('\n\n') if p.strip()]
    html_parts = [f'<p>{_esc(p)}</p>' for p in paras]
    return '\n'.join(html_parts), text, title


# ─── 统一入口 ────────────────────────────────────────────────

def _extract_document(raw: bytes, ext: str) -> Tuple[str, str, str]:
    """返回 (body_html, plain_text, title)"""
    if ext == 'pdf':
        has_gemini = bool(os.environ.get('GEMINI_API_KEY', ''))
        if has_gemini:
            return _extract_pdf_vision(raw)
        logger.warning('GEMINI_API_KEY 未设置，降级为字体分析提取')
        return _extract_pdf_dict(raw)
    elif ext == 'docx':
        return _extract_docx(raw)
    elif ext == 'pptx':
        return _extract_pptx(raw)
    elif ext in ('txt', 'md'):
        return _extract_text_file(raw)
    else:
        raise ValueError(f'不支持的格式：.{ext}')


# ─── 语言检测 ────────────────────────────────────────────────

def _detect_english(text: str) -> bool:
    sample = text[:2000]
    if not sample:
        return False
    non_ascii = sum(1 for c in sample if ord(c) > 127)
    return (non_ascii / len(sample)) < 0.4


# ─── HTML 组装 ───────────────────────────────────────────────

def _paragraphs_to_html(text: str) -> str:
    paras = [p.strip() for p in text.split('\n\n') if p.strip()]
    return ''.join(f'<p>{_esc(p)}</p>' for p in paras)


def _bilingual_html(text: str) -> str:
    raw_paras = [p.strip() for p in text.split('\n\n') if p.strip()]
    if not raw_paras:
        return ''
    fake_paras = [{'start': i, 'texts': [p]} for i, p in enumerate(raw_paras)]
    try:
        translations = _translate_paragraphs(fake_paras)
    except Exception as e:
        logger.warning(f'批量翻译失败，降级为原文: {e}')
        translations = [''] * len(raw_paras)
    parts = []
    for orig, zh in zip(raw_paras, translations):
        parts.append(f'<p lang="en">{_esc(orig)}</p>')
        if zh:
            parts.append(f'<p lang="zh">{_esc(zh)}</p>')
    return ''.join(parts)


def _build_summary_html(text: str) -> str:
    prompt = (
        '请为以下文档生成中文摘要与核心要点。'
        '返回 JSON 格式：{"summary": "一段话摘要", "key_points": ["要点1", "要点2", ...]}\n\n'
        f'文档内容（前 6000 字）：\n{text[:6000]}'
    )
    try:
        data = _llm_json(prompt, timeout=90)
        summary = data.get('summary', '')
        key_points = data.get('key_points', [])
    except Exception as e:
        logger.warning(f'摘要生成失败: {e}')
        return '<p>（摘要生成失败）</p>'
    points_html = ''.join(f'<li>{_esc(p)}</li>' for p in key_points) if key_points else ''
    return f'<p>{_esc(summary)}</p>' + (f'<ul>{points_html}</ul>' if points_html else '')


def _build_document_html(plain_text: str, is_english: bool, body_html: Optional[str] = None) -> str:
    summary_html = _build_summary_html(plain_text)
    if body_html:
        final_body = body_html
    elif is_english:
        final_body = _bilingual_html(plain_text)
    else:
        final_body = _paragraphs_to_html(plain_text)

    return (
        '<div class="doc-summary">'
        '<h2>摘要与核心要点</h2>'
        + summary_html
        + '</div>'
        '<hr>'
        '<div class="doc-content">'
        + final_body
        + '</div>'
    )


# ─── 入口 ────────────────────────────────────────────────────

async def handle_document(article_id: str) -> Dict[str, Any]:
    media_storage_path = os.environ.get('MEDIA_STORAGE_PATH', '')
    if not media_storage_path:
        return {'article_id': article_id, 'status': 'error',
                'error_message': 'MEDIA_STORAGE_PATH 未设置'}

    ext = media_storage_path.rsplit('.', 1)[-1].lower()

    if ext == 'doc':
        return {'article_id': article_id, 'status': 'error',
                'error_message': '暂不支持 .doc 格式，请转换为 .docx 后重新发送。'}
    if ext == 'ppt':
        return {'article_id': article_id, 'status': 'error',
                'error_message': '暂不支持 .ppt 格式，请转换为 .pptx 后重新发送。'}

    logger.info(f'下载文档：{media_storage_path}')
    raw_bytes = _download_from_storage(media_storage_path)

    logger.info(f'提取内容：ext={ext}，大小={len(raw_bytes)} bytes')
    body_html, plain_text, title = _extract_document(raw_bytes, ext)

    if not plain_text.strip():
        return {'article_id': article_id, 'status': 'error',
                'error_message': '文档内容为空，无法处理。'}

    is_english = _detect_english(plain_text)
    logger.info(f'语言：{"英文" if is_english else "中文"}，body_html={len(body_html)} chars')

    content_html = _build_document_html(plain_text, is_english, body_html=body_html)

    return {
        'article_id': article_id,
        'status': 'success',
        'title': title,
        'content_md': content_html,
        'content_text': plain_text[:4000],
        'video_storage_path': media_storage_path,
    }
